import streamlit as st
import os
import zipfile
from pptx import Presentation
from io import BytesIO
from pptx.shapes.autoshape import AutoShapeType

def split_ppt(uploaded_file, output_folder):
    """PPT 파일을 슬라이드별로 분할하고 압축합니다."""

    prs = Presentation(uploaded_file)
    file_list = []

    for i, slide in enumerate(prs.slides):
        new_prs = Presentation()
        new_slide = new_prs.slides.add_slide(new_prs.slide_layouts[0]) # 제목 슬라이드 레이아웃 사용 (원하는 레이아웃 선택 가능)
        
        for shape in slide.shapes:
            try:
                el = shape.element
                newel = new_slide.shapes.add_shape(el.prstGeom, el.x, el.y, el.cx, el.cy).element
                newel.sp.xml = el.sp.xml
            except KeyError as e:
                st.warning(f"슬라이드 {i+1}의 도형을 복사하는 중 오류 발생: {e}. 해당 도형은 건너뜁니다.")
            except Exception as e:
                st.error(f"예상치 못한 오류 발생: {e}")

        output_filename = f"slide_{i + 1}.pptx"
        output_path = os.path.join(output_folder, output_filename)
        new_prs.save(output_path)
        file_list.append(output_path)

    return file_list

def create_zip(file_list):
    """주어진 파일 목록을 zip 파일로 압축합니다."""

    zip_buffer = BytesIO()  # 메모리 버퍼에 zip 파일 생성

    with zipfile.ZipFile(zip_buffer, 'w') as zipf:
        for file_path in file_list:
            zipf.write(file_path, os.path.basename(file_path))

    zip_buffer.seek(0) # 버퍼의 시작 위치로 이동
    return zip_buffer


def main():
    st.title("PPT 분할 프로그램")

    uploaded_file = st.file_uploader("PPT 파일을 업로드하세요", type=["pptx"])

    if uploaded_file is not None:
        if st.button("실행"):
            temp_dir = "temp_ppt_split" #임시 폴더 생성
            os.makedirs(temp_dir, exist_ok=True)

            with st.spinner("PPT 분할 및 압축 중..."):
                file_list = split_ppt(uploaded_file, temp_dir)
                zip_buffer = create_zip(file_list)
                st.success("PPT 분할 및 압축 완료!")

            st.download_button(
                label="분할된 PPT 다운로드 (ZIP)",
                data=zip_buffer,
                file_name="split_ppt.zip",
                mime="application/zip",
            )

            # 임시 폴더 및 파일 삭제
            for file_path in file_list:
                os.remove(file_path)
            os.rmdir(temp_dir)

if __name__ == "__main__":
    main()